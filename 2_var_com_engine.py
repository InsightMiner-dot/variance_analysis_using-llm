import streamlit as st
import pandas as pd
import os
import re
from dotenv import load_dotenv
from typing import TypedDict, List, Dict, Any
from langgraph.graph import StateGraph, START, END
from langchain_openai import AzureChatOpenAI
from langchain_core.messages import SystemMessage, HumanMessage
from pydantic import BaseModel, Field
import io

from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from pptx import Presentation

# Load environment variables from .env
load_dotenv()

# ==========================================
# 1. PYDANTIC SCHEMAS (STRUCTURED OUTPUT)
# ==========================================
class FinalDriver(BaseModel):
    name: str = Field(description="Name of the variance driver (Plain text, no markdown)")
    variance: str = Field(description="Variance amount (e.g., '5.20M')")

class CategorySummary(BaseModel):
    primary_category: str = Field(description="Name of the primary category (Plain text)")
    drivers: List[FinalDriver] = Field(description="Top reasons/drivers from the final level")

class VarianceReport(BaseModel):
    ui_markdown_report: str = Field(description="The complete executive summary formatted as a Markdown string for the web UI.")
    overall_conclusion: str = Field(description="1-2 sentence overall executive conclusion (Plain text, no markdown)")
    category_breakdowns: List[CategorySummary] = Field(description="Breakdown of drivers per primary category")

# ==========================================
# 2. LANGGRAPH STATE DEFINITION
# ==========================================
class AgentState(TypedDict):
    df: pd.DataFrame
    hierarchy_cols: List[str]
    has_variance_col: bool
    variance_col: str
    base_scenario: str
    compare_scenario: str
    path_trace: List[str]
    final_summary: Any  

# ==========================================
# 3. LANGGRAPH NODES (WITH CLEAN MARKDOWN)
# ==========================================
def calculate_variance_node(state: AgentState) -> Dict[str, Any]:
    """Branched Drill-Down: Now outputs perfectly formatted Markdown for the UI."""
    df = state["df"].copy()
    hierarchy = state.get("hierarchy_cols", [])
    has_var = state.get("has_variance_col", True)
    
    path_trace = []
    
    if not hierarchy:
        return {"path_trace": ["⚠️ **Error:** No hierarchy columns selected."]}
        
    if has_var:
        target_col = state.get("variance_col")
        if target_col not in df.columns:
            return {"path_trace": [f"⚠️ **Error:** Target column '{target_col}' not found."]}
    else:
        base_col = state.get("base_scenario")
        comp_col = state.get("compare_scenario")
        if base_col not in df.columns or comp_col not in df.columns:
             return {"path_trace": ["⚠️ **Error:** Scenario columns not found."]}
        
        target_col = "Calculated_Variance"
        df[target_col] = df[base_col] - df[comp_col]

    total_variance = df[target_col].fillna(0).sum()
    path_trace.append(f"### 📈 Overall Total Variance: `{total_variance / 1e6:,.2f}M`\n---")

    first_level = hierarchy[0]
    first_level_grouped = df.groupby(first_level)[target_col].sum()
    top_primary_categories = first_level_grouped.reindex(first_level_grouped.abs().sort_values(ascending=False).index).head(5)

    for primary_cat, primary_val in top_primary_categories.items():
        path_trace.append(f"#### 🏢 Primary Category: **{primary_cat}** *(Total: {primary_val / 1e6:,.2f}M)*")
        
        current_df = df[df[first_level] == primary_cat]
        
        if len(hierarchy) == 1:
             path_trace.append(f"**FINAL LEVEL: Top 5 Reasons in '{first_level}'**")
             for idx, (name, val) in enumerate(top_primary_categories.items(), 1):
                 path_trace.append(f"- **{name}**: `{val / 1e6:,.2f}M`")
             path_trace.append("\n---\n")
             continue

        for i, level in enumerate(hierarchy[1:]):
            is_last_level = (i == len(hierarchy[1:]) - 1)
            
            if current_df.empty or level not in current_df.columns:
                break
                
            grouped = current_df.groupby(level)[target_col].sum()
            if grouped.empty or grouped.isna().all():
                break
                
            top_5 = grouped.reindex(grouped.abs().sort_values(ascending=False).index).head(5)
            
            if is_last_level:
                path_trace.append(f"**🎯 FINAL LEVEL: Top Reasons in '{level}'**")
            else:
                path_trace.append(f"**🔸 Top Drivers in '{level}'**")
                
            for idx, (name, val) in enumerate(top_5.items(), 1):
                path_trace.append(f"- {name}: `{val / 1e6:,.2f}M`")
                
            if is_last_level:
                break 
                
            max_driver = top_5.index[0]
            if pd.isna(max_driver):
                break
                
            path_trace.append(f"\n> *↳ Drilling down into '{max_driver}'...*\n")
            current_df = current_df[current_df[level] == max_driver]
            
        path_trace.append("\n---\n") 
        
    return {"path_trace": path_trace}

def synthesize_insight_node(state: AgentState) -> Dict[str, Any]:
    if state["path_trace"] and "⚠️ **Error:" in state["path_trace"][0]:
        return {"final_summary": "Analysis aborted due to invalid data configuration."}
    
    llm = AzureChatOpenAI(
        azure_endpoint=os.getenv("AZURE_OPENAI_ENDPOINT"),
        api_key=os.getenv("AZURE_OPENAI_KEY"),
        azure_deployment=os.getenv("AZURE_OPENAI_DEPLOYMENT"),
        api_version=os.getenv("AZURE_OPENAI_API_VERSION")
    )

    structured_llm = llm.with_structured_output(VarianceReport)

    system_prompt = (
        "You are a strict, professional financial data analyst. You are reviewing a branched variance "
        "analysis trace. All values are in Millions (M).\n\n"
        "Populate the JSON schema with pure text (no markdown asterisks or hashes in the drivers/categories).\n"
        "HOWEVER, for the 'ui_markdown_report' field, generate a highly readable Markdown string exactly as follows:\n"
        "1. A brief 1-2 sentence overall conclusion.\n"
        "2. A bulleted breakdown for each 'Primary Category' analyzed, listing the Top 5 reasons/drivers."
    )
    
    trace_text = "\n".join(state["path_trace"])
    messages = [
        SystemMessage(content=system_prompt),
        HumanMessage(content=f"Raw Data Trace:\n{trace_text}")
    ]
    
    response = structured_llm.invoke(messages)
    return {"final_summary": response.model_dump()}

# ==========================================
# 4. REPORT GENERATION UTILITIES
# ==========================================
def clean_markdown(text: str) -> str:
    """Strips Markdown symbols so the Excel file remains clean and professional."""
    text = re.sub(r'[*#_`]', '', text)
    text = text.replace('>', '')
    text = text.replace('↳', '->')
    return text.strip()

def generate_clean_excel(summary_data: dict, trace: List[str]) -> bytes:
    output = io.BytesIO()
    
    with pd.ExcelWriter(output, engine='openpyxl') as writer:
        summary_rows = []
        summary_rows.append({"Category": "OVERALL CONCLUSION", "Driver": summary_data["overall_conclusion"], "Amount": ""})
        summary_rows.append({"Category": "", "Driver": "", "Amount": ""}) 
        
        for cat in summary_data["category_breakdowns"]:
            summary_rows.append({"Category": cat['primary_category'], "Driver": "", "Amount": ""})
            for driver in cat["drivers"]:
                summary_rows.append({"Category": "", "Driver": driver["name"], "Amount": driver["variance"]})
        
        summary_df = pd.DataFrame(summary_rows)
        summary_df.to_excel(writer, sheet_name="Executive Summary", index=False)
        
        clean_trace = [clean_markdown(line) for line in trace if line.strip() != "---"]
        trace_df = pd.DataFrame({"Calculation Trace": clean_trace})
        trace_df.to_excel(writer, sheet_name="Variance Details", index=False)
        
        workbook = writer.book
        header_fill = PatternFill(start_color="1F4E78", end_color="1F4E78", fill_type="solid")
        header_font = Font(color="FFFFFF", bold=True, size=12)
        bold_font = Font(bold=True)
        cell_alignment = Alignment(wrap_text=True, vertical="top", horizontal="left")

        ws1 = workbook["Executive Summary"]
        for col in range(1, 4):
            ws1.cell(row=1, column=col).fill = header_fill
            ws1.cell(row=1, column=col).font = header_font
            ws1.cell(row=1, column=col).alignment = Alignment(horizontal="center", vertical="center")
            
        for row in range(2, ws1.max_row + 1):
            for col in range(1, 4):
                cell = ws1.cell(row=row, column=col)
                cell.alignment = cell_alignment
                if col == 1 and cell.value and cell.value != "OVERALL CONCLUSION":
                    cell.font = bold_font
                
        ws1.column_dimensions['A'].width = 35
        ws1.column_dimensions['B'].width = 60
        ws1.column_dimensions['C'].width = 15

        ws2 = workbook["Variance Details"]
        ws2.cell(row=1, column=1).fill = header_fill
        ws2.cell(row=1, column=1).font = header_font
        for row in range(2, ws2.max_row + 1):
            ws2.cell(row=row, column=1).alignment = cell_alignment
        ws2.column_dimensions['A'].width = 110
            
    return output.getvalue()

def generate_ppt_report(summary_data: dict) -> bytes:
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Root Cause Variance Analysis"
    slide.placeholders[1].text = "Automated Executive Briefing"
    
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Executive Conclusion"
    slide.shapes.placeholders[1].text_frame.text = summary_data["overall_conclusion"]
    
    for cat in summary_data["category_breakdowns"]:
        slide = prs.slides.add_slide(bullet_slide_layout)
        slide.shapes.title.text = f"Key Drivers: {cat['primary_category']}"
        tf = slide.shapes.placeholders[1].text_frame
        for driver in cat["drivers"]:
            p = tf.add_paragraph()
            p.text = f"{driver['name']}: {driver['variance']}"
            p.level = 0
            
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()

# ==========================================
# 5. GRAPH COMPILATION
# ==========================================
def build_graph():
    workflow = StateGraph(AgentState)
    workflow.add_node("calculate_variance", calculate_variance_node)
    workflow.add_node("synthesize_insight", synthesize_insight_node)
    workflow.add_edge(START, "calculate_variance")
    workflow.add_edge("calculate_variance", "synthesize_insight")
    workflow.add_edge("synthesize_insight", END)
    return workflow.compile()

# ==========================================
# 6. STREAMLIT UI & SESSION STATE MANAGEMENT
# ==========================================
st.set_page_config(page_title="Branched Variance Analyzer", layout="wide")

# Manage App Reset State
if "uploader_key" not in st.session_state:
    st.session_state.uploader_key = 0
if "analysis_result" not in st.session_state:
    st.session_state.analysis_result = None

def reset_app():
    """Clears the file uploader and analysis data to start fresh."""
    st.session_state.uploader_key += 1
    st.session_state.analysis_result = None

st.title("Branched Root Cause Analyzer")
st.write("Upload your dataset to generate a structured executive summary linking your primary categories directly to their root causes.")

with st.sidebar:
    st.header("1. Upload Data")
    # Dynamic key allows us to completely reset the file uploader without a hard page refresh
    uploaded_file = st.file_uploader("Upload CSV/Excel", type=["csv", "xlsx"], key=f"uploader_{st.session_state.uploader_key}")

    if uploaded_file is not None:
        try:
            if uploaded_file.name.endswith('.xlsx'):
                file_buffer = io.BytesIO(uploaded_file.getvalue())
                xls = pd.ExcelFile(file_buffer)
                sheet_name = st.selectbox("Select Sheet", xls.sheet_names)
                df = pd.read_excel(file_buffer, sheet_name=sheet_name)
            else:
                df = pd.read_csv(uploaded_file)
                
            st.header("2. Configure Metrics")
            num_cols = df.select_dtypes(include=['number']).columns.tolist()
            has_variance_col = st.checkbox("Variance column already present?", value=True)
            
            variance_col = ""
            base_scenario = ""
            compare_scenario = ""
            
            if has_variance_col:
                default_var = next((c for c in num_cols if 'var' in c.lower()), num_cols[0] if num_cols else None)
                idx = num_cols.index(default_var) if default_var in num_cols else 0
                variance_col = st.selectbox("Select Variance Column", num_cols, index=idx)
            else:
                st.caption("Calculate: (Base - Compare)")
                base_scenario = st.selectbox("Select Base Scenario", num_cols)
                compare_scenario = st.selectbox("Select Compare Scenario", num_cols)

            st.header("3. Select Hierarchy")
            cat_cols = df.select_dtypes(include=['object', 'string', 'category']).columns.tolist()
            all_cols = df.columns.tolist()
            
            hierarchy = st.multiselect("Hierarchy Flow", options=all_cols, default=cat_cols)

            if st.button("Generate Commentary", type="primary", use_container_width=True):
                if not hierarchy:
                    st.warning("Please select at least one column for the hierarchy.")
                elif has_variance_col and not variance_col:
                    st.warning("Please select a variance column.")
                elif not has_variance_col and (not base_scenario or not compare_scenario):
                    st.warning("Please select both Base and Compare scenarios.")
                else:
                    with st.spinner("Executing branched drill-down analysis..."):
                        app_graph = build_graph()
                        inputs = {
                            "df": df,
                            "hierarchy_cols": hierarchy,
                            "has_variance_col": has_variance_col,
                            "variance_col": variance_col,
                            "base_scenario": base_scenario,
                            "compare_scenario": compare_scenario
                        }
                        # Save result to session state so it persists during downloads
                        st.session_state.analysis_result = app_graph.invoke(inputs)

        except Exception as e:
            st.error(f"Error loading file: {e}")

# If we have an active analysis stored in session state, display it
if st.session_state.analysis_result is not None:
    result = st.session_state.analysis_result
    
    st.markdown("---")
    st.markdown("### 🤖 Executive Summary")
    
    summary_data = result["final_summary"]
    
    if isinstance(summary_data, str):
        st.error(summary_data)
    else:
        st.success(summary_data["ui_markdown_report"])
    
    st.markdown("### 🧮 Branched Drill-Down Trace")
    
    # Renders the trace using clean Markdown UI
    for step in result["path_trace"]:
        st.markdown(step)
            
    st.markdown("---")
    if not isinstance(summary_data, str):
        st.write("### 📥 Export Reports")
        
        col1, col2, col3 = st.columns([1, 1, 1])
        
        with col1:
            excel_data = generate_clean_excel(summary_data, result["path_trace"])
            st.download_button(
                label="📊 Download Excel",
                data=excel_data,
                file_name="Variance_Executive_Report.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                use_container_width=True
            )
            
        with col2:
            ppt_data = generate_ppt_report(summary_data)
            st.download_button(
                label="🖥️ Download PowerPoint",
                data=ppt_data,
                file_name="Variance_Executive_Briefing.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True
            )
            
        with col3:
            # Dedicated reset button - clears the UI and Uploader instantly
            st.button("🔄 Start New Analysis", on_click=reset_app, use_container_width=True, type="secondary")
